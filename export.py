"""
crop_tool - 统一导出模块

将图片按顺序导出为 PPT / PDF / Word。
所有导出功能共享相同的图片扫描和排序逻辑。
每个子文件夹生成一个独立的文件（{文件夹名}.pptx/pdf/docx）。

可独立运行：python export.py <input_dir> <output_dir> [ppt|pdf|word|all]
"""

import sys
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Inches
from docx import Document
from docx.enum.section import WD_ORIENT

import config
from utils import scan_speaker_dirs


def export_ppt(input_dir, output_dir, ratio=None):
    """
    导出为 PPT（每个子文件夹生成一个 .pptx 文件）

    参数:
        input_dir: 输入目录
        output_dir: 输出目录
        ratio: 视觉比例矫正（默认 config.TARGET_VISUAL_RATIO）
               16/9 = 铺满整页，4/3 = 居中留白矫正
    """
    if ratio is None:
        ratio = config.TARGET_VISUAL_RATIO

    output_path = Path(output_dir)
    output_path.mkdir(parents=True, exist_ok=True)

    speaker_dirs = scan_speaker_dirs(input_dir)
    if not speaker_dirs:
        print("  [错误] 没有找到任何图片，请检查输入目录。")
        return

    print(f"  共找到 {len(speaker_dirs)} 个文件夹\n")

    slide_w = config.SLIDE_WIDTH_INCH
    slide_h = config.SLIDE_HEIGHT_INCH
    slide_ratio = 16 / 9

    for speaker_name, images in speaker_dirs:
        print(f"  [{speaker_name}] {len(images)} 张图片")

        prs = Presentation()
        prs.slide_width = Inches(slide_w)
        prs.slide_height = Inches(slide_h)

        for idx, img_path in enumerate(images, 1):
            print(f"    {idx}/{len(images)}  {img_path.name}")
            try:
                blank_layout = prs.slide_layouts[6]
                slide = prs.slides.add_slide(blank_layout)

                if abs(ratio - 16/9) < 0.01:
                    slide.shapes.add_picture(str(img_path), 0, 0,
                                             width=Inches(slide_w),
                                             height=Inches(slide_h))
                else:
                    if ratio > slide_ratio:
                        display_width = Inches(slide_w)
                        display_height = Inches(slide_w / ratio)
                        left = 0
                        top = Inches((slide_h - slide_w / ratio) / 2)
                    else:
                        display_height = Inches(slide_h)
                        display_width = Inches(slide_h * ratio)
                        top = 0
                        left = Inches((slide_w - slide_h * ratio) / 2)

                    slide.shapes.add_picture(str(img_path), int(left), int(top),
                                             width=int(display_width),
                                             height=int(display_height))
            except Exception as e:
                print(f"    [失败] {e}")

        ppt_file = output_path / f"{speaker_name}.pptx"
        prs.save(str(ppt_file))
        print(f"    [完成] -> {ppt_file}\n")

    print("  [完成] PPT 导出完毕！")


def export_pdf(input_dir, output_dir):
    """
    导出为 PDF（每个子文件夹生成一个 .pdf 文件）

    使用 Pillow 的 Image.save(save_all=True) 功能，无需额外依赖。
    """
    output_path = Path(output_dir)
    output_path.mkdir(parents=True, exist_ok=True)

    speaker_dirs = scan_speaker_dirs(input_dir)
    if not speaker_dirs:
        print("  [错误] 没有找到任何图片，请检查输入目录。")
        return

    print(f"  共找到 {len(speaker_dirs)} 个文件夹\n")

    for speaker_name, images in speaker_dirs:
        print(f"  [{speaker_name}] {len(images)} 张图片")

        pil_images = []
        for idx, img_path in enumerate(images, 1):
            print(f"    {idx}/{len(images)}  {img_path.name}")
            try:
                img = Image.open(str(img_path)).convert("RGB")
                pil_images.append(img)
            except Exception as e:
                print(f"    [失败] {e}")

        if not pil_images:
            print(f"    [跳过] 无可用图片\n")
            continue

        pdf_file = output_path / f"{speaker_name}.pdf"
        pil_images[0].save(str(pdf_file), "PDF", save_all=True,
                           append_images=pil_images[1:])
        print(f"    [完成] -> {pdf_file}\n")

    print("  [完成] PDF 导出完毕！")


def export_word(input_dir, output_dir):
    """
    导出为 Word（每个子文件夹生成一个 .docx 文件）

    使用 python-docx，每张图片占一页。
    """
    output_path = Path(output_dir)
    output_path.mkdir(parents=True, exist_ok=True)

    speaker_dirs = scan_speaker_dirs(input_dir)
    if not speaker_dirs:
        print("  [错误] 没有找到任何图片，请检查输入目录。")
        return

    print(f"  共找到 {len(speaker_dirs)} 个文件夹\n")

    for speaker_name, images in speaker_dirs:
        print(f"  [{speaker_name}] {len(images)} 张图片")

        doc = Document()
        section = doc.sections[0]
        section.orientation = WD_ORIENT.LANDSCAPE
        section.page_width, section.page_height = section.page_height, section.page_width

        for idx, img_path in enumerate(images, 1):
            print(f"    {idx}/{len(images)}  {img_path.name}")
            try:
                pic_width = section.page_width - section.left_margin - section.right_margin
                doc.add_picture(str(img_path), width=pic_width)
                if idx < len(images):
                    doc.add_page_break()
            except Exception as e:
                print(f"    [失败] {e}")

        docx_file = output_path / f"{speaker_name}.docx"
        doc.save(str(docx_file))
        print(f"    [完成] -> {docx_file}\n")

    print("  [完成] Word 导出完毕！")


if __name__ == "__main__":
    if len(sys.argv) >= 4:
        input_dir = sys.argv[1]
        output_dir = sys.argv[2]
        fmt = sys.argv[3]
    else:
        print("用法: python export.py <input_dir> <output_dir> [ppt|pdf|word|all]")
        sys.exit(1)

    print("=" * 50)
    print("  图片导出工具 (PPT / PDF / Word)")
    print("=" * 50)
    print(f"  输入目录: {input_dir}")
    print(f"  输出目录: {output_dir}")
    print("=" * 50)

    if fmt in ("ppt", "all"):
        print("\n--- 导出 PPT ---")
        export_ppt(input_dir, output_dir)
    if fmt in ("pdf", "all"):
        print("\n--- 导出 PDF ---")
        export_pdf(input_dir, output_dir)
    if fmt in ("word", "all"):
        print("\n--- 导出 Word ---")
        export_word(input_dir, output_dir)
